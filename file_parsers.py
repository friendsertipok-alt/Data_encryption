import io
from docx import Document
import openpyxl
from pptx import Presentation

SENSITIVE_HEADERS = ["зарплата", "оклад", "премия", "доход", "бонус", "выплата", "сумма", "стоимость", "бюджет"]

def is_sensitive_header(header_text: str) -> bool:
    if not header_text: return False
    text = str(header_text).lower()
    for word in SENSITIVE_HEADERS:
        if word in text:
            return True
    return False

def anonymize_docx(file_bytes: bytes, anonymize_func) -> bytes:
    """Парсит Word документ, заменяет текст с сохранением стилей и возвращает байты."""
    doc = Document(io.BytesIO(file_bytes))
    
    # Замена в параграфах
    for para in doc.paragraphs:
        if para.text.strip():
            # Заменяем текст параграфа (в идеале нужно идти по run'ам для сохранения стилей каждого слова,
            # но для начала заменяем текст целиком)
            anonymized = anonymize_func(para.text)
            if anonymized != para.text:
                para.text = anonymized
                
    # Замена в таблицах
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    anonymized = anonymize_func(cell.text)
                    if anonymized != cell.text:
                        cell.text = anonymized
                        
    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()

def deanonymize_docx(file_bytes: bytes, deanonymize_func) -> bytes:
    """Восстанавливает текст в Word документе."""
    doc = Document(io.BytesIO(file_bytes))
    
    for para in doc.paragraphs:
        if para.text.strip():
            restored = deanonymize_func(para.text)
            if restored != para.text:
                para.text = restored
                
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    restored = deanonymize_func(cell.text)
                    if restored != cell.text:
                        cell.text = restored
                        
    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()

def anonymize_xlsx(file_bytes: bytes, anonymize_func, context_anonymize_func=None) -> bytes:
    """Парсит Excel документ, заменяет текст, учитывая заголовки столбцов."""
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    
    for sheet in wb.worksheets:
        sensitive_cols = set()
        
        # Определяем заголовки из первой строки
        for row in sheet.iter_rows(min_row=1, max_row=1):
            for col_idx, cell in enumerate(row):
                if cell.value and is_sensitive_header(str(cell.value)):
                    sensitive_cols.add(col_idx)
                    
        for row_idx, row in enumerate(sheet.iter_rows()):
            for col_idx, cell in enumerate(row):
                if cell.value is not None and isinstance(cell.value, (str, int, float)):
                    original_str = str(cell.value)
                    
                    # Если это чувствительная колонка (не заголовок) и в ней есть цифры
                    if row_idx > 0 and col_idx in sensitive_cols and context_anonymize_func and any(c.isdigit() for c in original_str):
                        anonymized = context_anonymize_func(original_str)
                    else:
                        anonymized = anonymize_func(original_str)
                        
                    if anonymized != original_str:
                        if isinstance(cell.value, int) and anonymized.isdigit():
                            cell.value = int(anonymized)
                        elif isinstance(cell.value, float):
                            try:
                                cell.value = float(anonymized)
                            except:
                                cell.value = anonymized
                        else:
                            cell.value = anonymized
                            
    output = io.BytesIO()
    wb.save(output)
    return output.getvalue()

def deanonymize_xlsx(file_bytes: bytes, deanonymize_func) -> bytes:
    """Восстанавливает текст в Excel документе."""
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
    
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is not None and isinstance(cell.value, (str, int, float)):
                    original_str = str(cell.value)
                    restored = deanonymize_func(original_str)
                    if restored != original_str:
                        if isinstance(cell.value, int) and restored.isdigit():
                            cell.value = int(restored)
                        elif isinstance(cell.value, float):
                            try:
                                cell.value = float(restored)
                            except:
                                cell.value = restored
                        else:
                            cell.value = restored
                        
    output = io.BytesIO()
    wb.save(output)
    return output.getvalue()

def _process_pptx_shapes(shapes, process_func):
    for shape in shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        processed = process_func(run.text)
                        if processed != run.text:
                            run.text = processed
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text.strip():
                                    processed = process_func(run.text)
                                    if processed != run.text:
                                        run.text = processed

def anonymize_pptx(file_bytes: bytes, anonymize_func) -> bytes:
    """Парсит PowerPoint документ, заменяет текст на слайдах."""
    prs = Presentation(io.BytesIO(file_bytes))
    for slide in prs.slides:
        _process_pptx_shapes(slide.shapes, anonymize_func)
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

def deanonymize_pptx(file_bytes: bytes, deanonymize_func) -> bytes:
    """Восстанавливает текст в PowerPoint документе."""
    prs = Presentation(io.BytesIO(file_bytes))
    for slide in prs.slides:
        _process_pptx_shapes(slide.shapes, deanonymize_func)
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
