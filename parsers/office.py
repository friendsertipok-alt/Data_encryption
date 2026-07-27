import io
from docx import Document
import openpyxl
from pptx import Presentation

SENSITIVE_HEADERS = ["зарплата", "оклад", "премия", "доход", "бонус", "выплата", "сумма", "стоимость", "бюджет", "salary", "bonus", "amount"]

def is_sensitive_header(header_text: str) -> bool:
    if not header_text: return False
    text = str(header_text).lower()
    for word in SENSITIVE_HEADERS:
        if word in text:
            return True
    return False

class OfficeParser:
    @staticmethod
    def _process_docx_paragraph(para, process_func):
        full_text = para.text
        if not full_text.strip():
            return
        processed = process_func(full_text)
        if processed == full_text:
            return
        if len(para.runs) == 1:
            para.runs[0].text = processed
        elif len(para.runs) > 1:
            para.runs[0].text = processed
            for run in para.runs[1:]:
                run.text = ""
        else:
            para.text = processed

    @staticmethod
    def _process_docx_runs(doc, process_func):
        """Обрабатывает текст через runs, сохраняя форматирование параграфов, таблиц и колонтитулов."""
        # Основной текст
        for para in doc.paragraphs:
            OfficeParser._process_docx_paragraph(para, process_func)

        # Таблицы
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        OfficeParser._process_docx_paragraph(para, process_func)

        # Колонтитулы (Headers & Footers)
        for section in doc.sections:
            if section.header:
                for para in section.header.paragraphs:
                    OfficeParser._process_docx_paragraph(para, process_func)
                for table in section.header.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                OfficeParser._process_docx_paragraph(para, process_func)
            if section.footer:
                for para in section.footer.paragraphs:
                    OfficeParser._process_docx_paragraph(para, process_func)
                for table in section.footer.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                OfficeParser._process_docx_paragraph(para, process_func)

    @staticmethod
    def anonymize_docx(file_bytes: bytes, anonymize_func) -> bytes:
        doc = Document(io.BytesIO(file_bytes))
        OfficeParser._process_docx_runs(doc, anonymize_func)
        output = io.BytesIO()
        doc.save(output)
        return output.getvalue()

    @staticmethod
    def deanonymize_docx(file_bytes: bytes, deanonymize_func) -> bytes:
        doc = Document(io.BytesIO(file_bytes))
        OfficeParser._process_docx_runs(doc, deanonymize_func)
        output = io.BytesIO()
        doc.save(output)
        return output.getvalue()

    @staticmethod
    def _load_workbook(file_bytes: bytes) -> openpyxl.Workbook:
        """Загружает файл Excel (.xlsx или .xls), автоматическая конвертация BIFF .xls с сохранением верстки."""
        if file_bytes.startswith(b'\xd0\xcf\x11\xe0'):
            import xlrd
            from openpyxl.utils import get_column_letter
            wb_xls = xlrd.open_workbook(file_contents=file_bytes, formatting_info=True)
            wb_xlsx = openpyxl.Workbook()
            wb_xlsx.remove(wb_xlsx.active)
            for name in wb_xls.sheet_names():
                sheet_xls = wb_xls.sheet_by_name(name)
                ws = wb_xlsx.create_sheet(title=name)
                # 1. Значения ячеек
                for r in range(sheet_xls.nrows):
                    for c in range(sheet_xls.ncols):
                        val = sheet_xls.cell_value(r, c)
                        ws.cell(row=r+1, column=c+1, value=val)
                # 2. Объединенные ячейки (Merged Cells) — ключевое для верстки!
                for crange in sheet_xls.merged_cells:
                    rlow, rhigh, clow, chigh = crange
                    ws.merge_cells(start_row=rlow+1, start_column=clow+1, end_row=rhigh, end_column=chigh)
                # 3. Ширины колонок
                if hasattr(sheet_xls, 'colinfo_map'):
                    for col_idx, colinfo in sheet_xls.colinfo_map.items():
                        col_letter = get_column_letter(col_idx + 1)
                        if colinfo.width:
                            ws.column_dimensions[col_letter].width = max(colinfo.width / 256.0, 10.0)
                # 4. Высоты строк
                if hasattr(sheet_xls, 'rowinfo_map'):
                    for row_idx, rowinfo in sheet_xls.rowinfo_map.items():
                        if rowinfo.height:
                            ws.row_dimensions[row_idx + 1].height = rowinfo.height / 20.0
            return wb_xlsx
        else:
            return openpyxl.load_workbook(io.BytesIO(file_bytes))

    @staticmethod
    def anonymize_xlsx(file_bytes: bytes, anonymize_func, context_anonymize_func=None) -> bytes:
        wb = OfficeParser._load_workbook(file_bytes)
        for sheet in wb.worksheets:
            sensitive_cols = set()
            for row in sheet.iter_rows(min_row=1, max_row=1):
                for col_idx, cell in enumerate(row):
                    if cell.value and is_sensitive_header(str(cell.value)):
                        sensitive_cols.add(col_idx)
            for row_idx, row in enumerate(sheet.iter_rows()):
                for col_idx, cell in enumerate(row):
                    if cell.value is None:
                        continue
                    if isinstance(cell.value, str) and cell.value.strip():
                        anonymized = anonymize_func(cell.value)
                        if anonymized != cell.value:
                            cell.value = anonymized
                    elif row_idx > 0 and col_idx in sensitive_cols and context_anonymize_func and isinstance(cell.value, (int, float)):
                        original_str = str(cell.value)
                        anonymized = context_anonymize_func(original_str)
                        if anonymized != original_str:
                            try:
                                cell.value = float(anonymized) if isinstance(cell.value, float) else int(anonymized)
                            except:
                                cell.value = anonymized
        output = io.BytesIO()
        wb.save(output)
        return output.getvalue()

    @staticmethod
    def deanonymize_xlsx(file_bytes: bytes, deanonymize_func) -> bytes:
        wb = OfficeParser._load_workbook(file_bytes)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None and isinstance(cell.value, str) and cell.value.strip():
                        restored = deanonymize_func(cell.value)
                        if restored != cell.value:
                            cell.value = restored
        output = io.BytesIO()
        wb.save(output)
        return output.getvalue()

    @staticmethod
    def _process_pptx_shapes(shapes, process_func):
        for shape in shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            processed = process_func(run.text)
                            if processed != run.text:
                                run.text = processed
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.text.strip():
                                        processed = process_func(run.text)
                                        if processed != run.text:
                                            run.text = processed

    @staticmethod
    def anonymize_pptx(file_bytes: bytes, anonymize_func) -> bytes:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            OfficeParser._process_pptx_shapes(slide.shapes, anonymize_func)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                for para in slide.notes_slide.notes_text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            processed = anonymize_func(run.text)
                            if processed != run.text:
                                run.text = processed
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    @staticmethod
    def deanonymize_pptx(file_bytes: bytes, deanonymize_func) -> bytes:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            OfficeParser._process_pptx_shapes(slide.shapes, deanonymize_func)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                for para in slide.notes_slide.notes_text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            processed = deanonymize_func(run.text)
                            if processed != run.text:
                                run.text = processed
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()
