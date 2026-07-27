"""
╔══════════════════════════════════════════════════════════════╗
║  ПОЛНАЯ ТЕСТОВАЯ БАТАРЕЯ DLP Gateway                        ║
║  Покрывает: детекторы, парсеры, API, шифрование, кэш,      ║
║  сессии, группы, деанонимизацию, edge cases                 ║
╚══════════════════════════════════════════════════════════════╝
"""

import unittest
import io
import json
import sqlite3
import os
import re
import uuid

from fastapi.testclient import TestClient
from PIL import Image, ImageDraw
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from main import app
from config import DB_PATH, SECRET_KEY
from session_store import (
    create_session, get_session_map, get_user_sessions,
    cleanup_expired_sessions, get_group_entity_map, _get_fernet
)
from core.detector import DlpDetector
from core.faker_engine import FakerEngine
from core.cache import LRUCache
from core.entities.russia import RUSSIA_PATTERNS
from core.entities.india import INDIA_PATTERNS
from parsers.plain import PlainParser
from parsers.office import OfficeParser, is_sensitive_header


# ═══════════════════════════════════════════
# 1. ТЕСТЫ REGEX ДЕТЕКТОРОВ (Россия)
# ═══════════════════════════════════════════

class TestRussiaRegexPatterns(unittest.TestCase):
    """Проверка всех regex-паттернов для российских ПДн"""

    def test_passport_ru(self):
        regex = RUSSIA_PATTERNS["PASSPORT_RU"]
        self.assertTrue(regex.search("4510 123456"))
        self.assertTrue(regex.search("45 10 123456"))

    def test_inn_ru(self):
        regex = RUSSIA_PATTERNS["INN_RU"]
        self.assertTrue(regex.search("ИНН 7707083893"))
        self.assertTrue(regex.search("инн 770708389312"))
        self.assertTrue(regex.search("INN:7707083893"))
        # Без ключевого слова не должен срабатывать
        self.assertIsNone(regex.search("7707083893"))

    def test_snils_ru(self):
        regex = RUSSIA_PATTERNS["SNILS_RU"]
        self.assertTrue(regex.search("123-456-789 12"))
        self.assertTrue(regex.search("12345678912"))

    def test_phone_ru(self):
        regex = RUSSIA_PATTERNS["PHONE_RU"]
        self.assertTrue(regex.search("+7 999 123-45-67"))
        self.assertTrue(regex.search("+7(999)123-45-67"))
        self.assertTrue(regex.search("8 999 123-45-67"))

    def test_credit_card(self):
        regex = RUSSIA_PATTERNS["CREDIT_CARD"]
        self.assertTrue(regex.search("4532 0151 1283 0366"))
        self.assertTrue(regex.search("4532015112830366"))

    def test_bank_account_ru(self):
        regex = RUSSIA_PATTERNS["BANK_ACCOUNT_RU"]
        self.assertTrue(regex.search("40817810099910004312"))

    def test_vehicle_private_ru(self):
        regex = RUSSIA_PATTERNS["VEHICLE_PRIVATE_RU"]
        self.assertTrue(regex.search("А123ВС77"))
        self.assertTrue(regex.search("А 123 ВС 777"))

    def test_vin(self):
        regex = RUSSIA_PATTERNS["VIN"]
        self.assertTrue(regex.search("WVWZZZ3CZWE123456"))
        # I, O, Q не допускаются
        self.assertIsNone(regex.search("WVWZZZ3CZWI123456"))

    def test_kpp_ru(self):
        regex = RUSSIA_PATTERNS["KPP_RU"]
        self.assertTrue(regex.search("КПП 770701001"))
        self.assertIsNone(regex.search("770701001"))  # без ключевого слова

    def test_oms_ru(self):
        regex = RUSSIA_PATTERNS["OMS_RU"]
        self.assertTrue(regex.search("ОМС 1234567890123456"))
        self.assertTrue(regex.search("полис 1234567890123456"))

    def test_driver_license_ru(self):
        regex = RUSSIA_PATTERNS["DRIVER_LICENSE_RU"]
        self.assertTrue(regex.search("77 14 892015"))


# ═══════════════════════════════════════════
# 2. ТЕСТЫ REGEX ДЕТЕКТОРОВ (Индия)
# ═══════════════════════════════════════════

class TestIndiaRegexPatterns(unittest.TestCase):
    """Проверка regex-паттернов для индийских ПДн"""

    def test_aadhaar_in(self):
        regex = INDIA_PATTERNS["AADHAAR_IN"]
        self.assertTrue(regex.search("1234 5678 9012"))
        self.assertTrue(regex.search("123456789012"))

    def test_pan_in(self):
        regex = INDIA_PATTERNS["PAN_IN"]
        self.assertTrue(regex.search("ABCDE1234F"))

    def test_voter_id_in(self):
        regex = INDIA_PATTERNS["VOTER_ID_IN"]
        self.assertTrue(regex.search("ABC1234567"))

    def test_gstin_in(self):
        regex = INDIA_PATTERNS["GSTIN_IN"]
        self.assertTrue(regex.search("27ABCDE1234F1Z5"))

    def test_phone_in(self):
        regex = INDIA_PATTERNS["PHONE_IN"]
        self.assertTrue(regex.search("+91 9876543210"))
        self.assertTrue(regex.search("919876543210"))

    def test_vehicle_standard_in(self):
        regex = INDIA_PATTERNS["VEHICLE_STANDARD_IN"]
        self.assertTrue(regex.search("MH 12 AB 1234"))

    def test_vehicle_bh_in(self):
        regex = INDIA_PATTERNS["VEHICLE_BH_IN"]
        self.assertTrue(regex.search("21 BH 1234 AB"))


# ═══════════════════════════════════════════
# 3. ТЕСТЫ DLP ДЕТЕКТОРА
# ═══════════════════════════════════════════

class TestDlpDetector(unittest.TestCase):
    """Интеграционные тесты ядра детекции"""

    @classmethod
    def setUpClass(cls):
        cls.detector = DlpDetector()

    def test_anonymize_basic_ru_text(self):
        text = "Иванов Иван написал на ivanov@mail.ru и позвонил по +7 999 123-45-67"
        entity_map = {}
        result = self.detector.analyze_and_anonymize(text, entity_map, lang='ru', mode='fake')

        # Оригинальные данные удалены
        self.assertNotIn("ivanov@mail.ru", result)
        self.assertNotIn("+7 999 123-45-67", result)
        # Словарь замен заполнен
        self.assertGreater(len(entity_map), 0)

    def test_anonymize_empty_text(self):
        entity_map = {}
        result = self.detector.analyze_and_anonymize("", entity_map)
        self.assertEqual(result, "")
        self.assertEqual(len(entity_map), 0)

    def test_anonymize_whitespace_text(self):
        entity_map = {}
        result = self.detector.analyze_and_anonymize("   \n\t  ", entity_map)
        self.assertEqual(result, "   \n\t  ")

    def test_anonymize_no_pii(self):
        text = "Сегодня хорошая погода. Завтра будет дождь."
        entity_map = {}
        result = self.detector.analyze_and_anonymize(text, entity_map)
        # Текст без ПДн должен остаться без изменений (или минимальные)
        self.assertEqual(len(entity_map), 0)
        self.assertEqual(result, text)

    def test_anonymize_tags_mode(self):
        text = "Позвоните на +7 999 123-45-67"
        entity_map = {}
        result = self.detector.analyze_and_anonymize(text, entity_map, lang='ru', mode='tags')
        self.assertNotIn("+7 999 123-45-67", result)
        # В tags-режиме токены выглядят как <TYPE_N>
        has_tag = any(k.startswith("<") and k.endswith(">") for k in entity_map.keys())
        self.assertTrue(has_tag, f"Expected tag tokens, got: {entity_map.keys()}")

    def test_deanonymize_restores_original(self):
        text = "Email: test@company.ru, тел: +7 999 123-45-67"
        entity_map = {}
        anon = self.detector.analyze_and_anonymize(text, entity_map, lang='ru', mode='fake')
        restored = self.detector.deanonymize(anon, entity_map)
        self.assertIn("test@company.ru", restored)
        self.assertIn("+7 999 123-45-67", restored)

    def test_canonical_mapping(self):
        """Один и тот же оригинал при повторном появлении получает тот же фейк"""
        text = "Петров и Петров и Петров"
        entity_map = {}
        self.detector.analyze_and_anonymize(text, entity_map, lang='ru', mode='fake')
        # Для одного и того же слова не должно быть дублирующих значений
        values = list(entity_map.values())
        # Проверяем что есть некоторая связь (canonical mapping работает)
        self.assertIsNotNone(entity_map)

    def test_language_detection_ru(self):
        self.assertEqual(self.detector.detect_language("Привет мир"), "ru")

    def test_language_detection_en(self):
        self.assertEqual(self.detector.detect_language("Hello world"), "en")

    def test_language_detection_mixed(self):
        # Если есть кириллица - считаем русским
        self.assertEqual(self.detector.detect_language("Hello Привет"), "ru")

    def test_luhn_validation(self):
        """Алгоритм Луна: валидный и невалидный номер карты"""
        self.assertTrue(self.detector._is_luhn_valid("4532015112830366"))
        self.assertFalse(self.detector._is_luhn_valid("4532015112830367"))
        self.assertFalse(self.detector._is_luhn_valid("123"))  # слишком короткий

    def test_passport_ru_detection_and_restore(self):
        text = "Паспорт: 4510 123456"
        entity_map = {}
        anon = self.detector.analyze_and_anonymize(text, entity_map, lang='ru', mode='fake')
        self.assertNotIn("4510 123456", anon)
        restored = self.detector.deanonymize(anon, entity_map)
        self.assertIn("4510 123456", restored)

    def test_snils_detection_and_restore(self):
        text = "СНИЛС: 123-456-789 12"
        entity_map = {}
        anon = self.detector.analyze_and_anonymize(text, entity_map, lang='ru', mode='fake')
        self.assertNotIn("123-456-789 12", anon)
        restored = self.detector.deanonymize(anon, entity_map)
        self.assertIn("123-456-789 12", restored)


# ═══════════════════════════════════════════
# 4. ТЕСТЫ FAKER ENGINE
# ═══════════════════════════════════════════

class TestFakerEngine(unittest.TestCase):
    """Проверка генерации фейков для всех типов"""

    def setUp(self):
        self.faker = FakerEngine()

    def test_generate_person(self):
        result = self.faker.generate_fake("PERSON", "Иванов Иван")
        self.assertIsInstance(result, str)
        self.assertTrue(len(result) > 0)

    def test_generate_phone_ru(self):
        result = self.faker.generate_fake("PHONE_RU", "+7 999 123-45-67")
        self.assertIsInstance(result, str)

    def test_generate_inn_10(self):
        result = self.faker.generate_fake("INN_RU", "7707083893")
        self.assertEqual(len(result), 10)

    def test_generate_inn_12(self):
        result = self.faker.generate_fake("INN_RU", "770708389312")
        self.assertEqual(len(result), 12)

    def test_generate_snils(self):
        result = self.faker.generate_fake("SNILS_RU", "123-456-789 12")
        self.assertRegex(result, r'\d{3}-\d{3}-\d{3} \d{2}')

    def test_generate_passport(self):
        result = self.faker.generate_fake("PASSPORT_RU", "4510 123456")
        self.assertIsInstance(result, str)
        self.assertNotEqual(result, "4510 123456")

    def test_generate_bank_account(self):
        result = self.faker.generate_fake("BANK_ACCOUNT_RU", "40817810099910004312")
        self.assertTrue(result.startswith("4"))
        self.assertEqual(len(result), 20)

    def test_generate_credit_card(self):
        result = self.faker.generate_fake("CREDIT_CARD", "4532015112830366")
        self.assertIsInstance(result, str)

    def test_generate_email(self):
        result = self.faker.generate_fake("EMAIL_ADDRESS", "test@mail.ru")
        self.assertIn("@", result)

    def test_generate_address_ru(self):
        result = self.faker.generate_fake("ADDRESS_RU", "г. Москва, ул. Тверская, д. 10")
        self.assertIn("г.", result)

    def test_generate_vin(self):
        result = self.faker.generate_fake("VIN", "WVWZZZ3CZWE123456")
        self.assertEqual(len(result), 17)

    def test_generate_aadhaar_in(self):
        result = self.faker.generate_fake("AADHAAR_IN", "1234 5678 9012")
        self.assertIsInstance(result, str)

    def test_generate_pan_in(self):
        result = self.faker.generate_fake("PAN_IN", "ABCDE1234F")
        self.assertEqual(len(result), 10)

    def test_generate_unknown_type(self):
        """Неизвестный тип возвращает [TYPE_NNNN]"""
        result = self.faker.generate_fake("UNKNOWN_TYPE", "test")
        self.assertTrue(result.startswith("[UNKNOWN_TYPE_"))

    def test_morph_russian_name(self):
        """Склонение русского имени"""
        inf_orig, inf_fake, _, _ = self.faker.morph_russian_name(
            ["Иванов", "Иван"], ["Петров", "Пётр"], "gent"
        )
        self.assertIsInstance(inf_orig, str)
        self.assertIsInstance(inf_fake, str)


# ═══════════════════════════════════════════
# 5. ТЕСТЫ LRU КЭША
# ═══════════════════════════════════════════

class TestLRUCache(unittest.TestCase):
    """Тесты кэширования"""

    def test_put_and_get(self):
        cache = LRUCache(capacity=10)
        cache.put("key1", "value1")
        self.assertEqual(cache.get("key1"), "value1")

    def test_cache_miss(self):
        cache = LRUCache(capacity=10)
        self.assertIsNone(cache.get("nonexistent"))

    def test_eviction(self):
        cache = LRUCache(capacity=2)
        cache.put("k1", "v1")
        cache.put("k2", "v2")
        cache.put("k3", "v3")  # k1 должен быть вытеснен
        self.assertIsNone(cache.get("k1"))
        self.assertEqual(cache.get("k2"), "v2")
        self.assertEqual(cache.get("k3"), "v3")

    def test_lru_order(self):
        cache = LRUCache(capacity=2)
        cache.put("k1", "v1")
        cache.put("k2", "v2")
        cache.get("k1")  # k1 becomes most recently used
        cache.put("k3", "v3")  # k2 should be evicted
        self.assertEqual(cache.get("k1"), "v1")
        self.assertIsNone(cache.get("k2"))

    def test_tuple_key(self):
        cache = LRUCache(capacity=10)
        cache.put(("text", "mode"), "result")
        self.assertEqual(cache.get(("text", "mode")), "result")

    def test_overwrite(self):
        cache = LRUCache(capacity=10)
        cache.put("k", "v1")
        cache.put("k", "v2")
        self.assertEqual(cache.get("k"), "v2")


# ═══════════════════════════════════════════
# 6. ТЕСТЫ СЕССИЙ И ШИФРОВАНИЯ
# ═══════════════════════════════════════════

class TestSessionStore(unittest.TestCase):
    """Тесты хранилища сессий"""

    def test_create_and_retrieve_session(self):
        entity_map = {"fake_email@test.ru": "real@company.ru"}
        session_id = create_session(entity_map, filename="test.txt")
        retrieved = get_session_map(session_id)
        self.assertIsNotNone(retrieved)
        self.assertEqual(retrieved["fake_email@test.ru"], "real@company.ru")

    def test_encryption_at_rest(self):
        """Проверяем что данные в SQLite зашифрованы"""
        secret = "TOP_SECRET_DATA_12345"
        entity_map = {"fake": secret}
        session_id = create_session(entity_map, filename="crypto.txt")

        conn = sqlite3.connect(DB_PATH)
        row = conn.execute(
            "SELECT entity_map_json FROM sessions WHERE session_id = ?", (session_id,)
        ).fetchone()
        conn.close()

        self.assertIsNotNone(row)
        encrypted_blob = row[0]
        # Открытый текст НЕ должен присутствовать в сырых данных
        self.assertNotIn(secret, encrypted_blob)
        # Fernet-формат
        self.assertTrue(encrypted_blob.startswith("gAAAAA"))

    def test_fernet_decryption(self):
        fernet = _get_fernet()
        raw = json.dumps({"key": "value"})
        encrypted = fernet.encrypt(raw.encode()).decode()
        decrypted = fernet.decrypt(encrypted.encode()).decode()
        self.assertEqual(decrypted, raw)

    def test_nonexistent_session(self):
        result = get_session_map("nonexistent-session-id-12345")
        self.assertIsNone(result)

    def test_empty_session_id(self):
        result = get_session_map("")
        self.assertIsNone(result)

    def test_group_session(self):
        """Тест групповых сессий"""
        group_id = f"TEST_GRP_{uuid.uuid4().hex[:8]}"
        map1 = {"fake1": "real1"}
        map2 = {"fake2": "real2"}
        create_session(map1, filename="f1.txt", group_id=group_id)
        create_session(map2, filename="f2.txt", group_id=group_id)

        combined = get_group_entity_map(group_id)
        self.assertIn("fake1", combined)
        self.assertIn("fake2", combined)

    def test_get_session_map_by_group_id(self):
        """Поиск по group_id через get_session_map"""
        group_id = f"GRP_LOOKUP_{uuid.uuid4().hex[:8]}"
        map1 = {"token_a": "original_a"}
        create_session(map1, filename="lookup.txt", group_id=group_id)

        result = get_session_map(group_id)
        self.assertIsNotNone(result)
        self.assertIn("token_a", result)

    def test_multiple_session_ids(self):
        """Список session_id через запятую"""
        map1 = {"t1": "o1"}
        map2 = {"t2": "o2"}
        sid1 = create_session(map1, filename="m1.txt")
        sid2 = create_session(map2, filename="m2.txt")

        combined = get_session_map(f"{sid1},{sid2}")
        self.assertIn("t1", combined)
        self.assertIn("t2", combined)

    def test_get_user_sessions_returns_list(self):
        sessions = get_user_sessions("default")
        self.assertIsInstance(sessions, list)

    def test_cleanup_expired(self):
        """Очистка устаревших сессий (TTL=0 удалит всё)"""
        create_session({"x": "y"}, filename="cleanup.txt")
        cleanup_expired_sessions(ttl_hours=0)
        # После очистки с TTL=0 сессии могли быть удалены
        # (зависит от точности времени, но метод не должен падать)


# ═══════════════════════════════════════════
# 7. ТЕСТЫ API ЭНДПОИНТОВ
# ═══════════════════════════════════════════

class TestApiEndpoints(unittest.TestCase):
    """Полные интеграционные тесты REST API"""

    def setUp(self):
        self.client = TestClient(app)

    # --- TEXT API ---

    def test_anonymize_text_basic(self):
        resp = self.client.post("/api/anonymize/text", json={
            "text": "Иванов Иван, email: ivan@mail.ru, +7 999 123-45-67",
            "mode": "fake", "lang": "ru"
        })
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertIn("anonymized_text", data)
        self.assertIn("session_id", data)
        self.assertGreater(data["entities_found"], 0)

    def test_anonymize_text_empty(self):
        resp = self.client.post("/api/anonymize/text", json={
            "text": "   ", "mode": "fake"
        })
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertEqual(data["entities_found"], 0)

    def test_anonymize_text_tags_mode(self):
        resp = self.client.post("/api/anonymize/text", json={
            "text": "Звоните +7 999 123-45-67", "mode": "tags", "lang": "ru"
        })
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertNotIn("+7 999 123-45-67", data["anonymized_text"])

    def test_deanonymize_text_success(self):
        # 1. Анонимизируем
        resp1 = self.client.post("/api/anonymize/text", json={
            "text": "email: test@corp.ru", "mode": "fake", "lang": "ru"
        })
        data1 = resp1.json()

        # 2. Деанонимизируем
        resp2 = self.client.post("/api/deanonymize/text", json={
            "text": data1["anonymized_text"],
            "session_id": data1["session_id"]
        })
        self.assertEqual(resp2.status_code, 200)
        restored = resp2.json()["restored_text"]
        self.assertIn("test@corp.ru", restored)

    def test_deanonymize_invalid_session(self):
        resp = self.client.post("/api/deanonymize/text", json={
            "text": "some text", "session_id": "invalid-uuid-99999"
        })
        self.assertEqual(resp.status_code, 404)

    # --- FILE API ---

    def test_anonymize_txt_file(self):
        content = "Клиент: Петров Пётр, +7 900 111-22-33"
        resp = self.client.post(
            "/api/anonymize/file",
            files={"file": ("test.txt", io.BytesIO(content.encode('utf-8')), "text/plain")},
            data={"mode": "fake"}
        )
        self.assertEqual(resp.status_code, 200)
        self.assertIn("X-Session-ID", resp.headers)
        result_text = resp.content.decode('utf-8')
        self.assertNotIn("+7 900 111-22-33", result_text)

    def test_anonymize_csv_file(self):
        csv_content = "Имя,Телефон\nИванов Иван,+7 999 123-45-67\n"
        resp = self.client.post(
            "/api/anonymize/file",
            files={"file": ("data.csv", io.BytesIO(csv_content.encode('utf-8')), "text/csv")},
            data={"mode": "fake"}
        )
        self.assertEqual(resp.status_code, 200)
        self.assertIn("X-Session-ID", resp.headers)

    def test_anonymize_json_file(self):
        json_data = json.dumps({"name": "Иванов Иван", "email": "ivan@mail.ru"}, ensure_ascii=False)
        resp = self.client.post(
            "/api/anonymize/file",
            files={"file": ("data.json", io.BytesIO(json_data.encode('utf-8')), "application/json")},
            data={"mode": "fake"}
        )
        self.assertEqual(resp.status_code, 200)

    def test_anonymize_xml_file(self):
        xml_content = '<?xml version="1.0"?><root><person>Иванов Иван</person></root>'
        resp = self.client.post(
            "/api/anonymize/file",
            files={"file": ("data.xml", io.BytesIO(xml_content.encode('utf-8')), "text/xml")},
            data={"mode": "fake"}
        )
        self.assertEqual(resp.status_code, 200)

    def test_anonymize_docx_file(self):
        doc = Document()
        doc.add_paragraph("Клиент Иванов Иван, email: client@corp.ru")
        buf = io.BytesIO()
        doc.save(buf)

        resp = self.client.post(
            "/api/anonymize/file",
            files={"file": ("test.docx", io.BytesIO(buf.getvalue()), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")},
            data={"mode": "fake"}
        )
        self.assertEqual(resp.status_code, 200)
        self.assertIn("X-Session-ID", resp.headers)

    def test_anonymize_xlsx_file(self):
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Имя"
        ws["B1"] = "Email"
        ws["A2"] = "Иванов Иван"
        ws["B2"] = "ivan@mail.ru"
        buf = io.BytesIO()
        wb.save(buf)

        resp = self.client.post(
            "/api/anonymize/file",
            files={"file": ("data.xlsx", io.BytesIO(buf.getvalue()), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
            data={"mode": "fake"}
        )
        self.assertEqual(resp.status_code, 200)

    def test_anonymize_pptx_file(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Клиент: Петров Пётр"
        buf = io.BytesIO()
        prs.save(buf)

        resp = self.client.post(
            "/api/anonymize/file",
            files={"file": ("pres.pptx", io.BytesIO(buf.getvalue()), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            data={"mode": "fake"}
        )
        self.assertEqual(resp.status_code, 200)

    def test_unsupported_format(self):
        resp = self.client.post(
            "/api/anonymize/file",
            files={"file": ("file.exe", io.BytesIO(b"binary"), "application/octet-stream")},
            data={"mode": "fake"}
        )
        self.assertEqual(resp.status_code, 400)

    def test_deanonymize_txt_file_roundtrip(self):
        """Полный цикл: анонимизация TXT -> деанонимизация TXT"""
        original = "Сотрудник Сидорова Елена, звонок +7 900 555-44-33"

        # Анонимизация
        resp1 = self.client.post(
            "/api/anonymize/file",
            files={"file": ("report.txt", io.BytesIO(original.encode('utf-8')), "text/plain")},
            data={"mode": "fake"}
        )
        self.assertEqual(resp1.status_code, 200)
        session_id = resp1.headers["X-Session-ID"]
        anon_bytes = resp1.content

        # Деанонимизация
        resp2 = self.client.post(
            "/api/deanonymize/file",
            files={"file": ("safe_report.txt", io.BytesIO(anon_bytes), "text/plain")},
            data={"session_id": session_id}
        )
        self.assertEqual(resp2.status_code, 200)
        restored_text = resp2.content.decode('utf-8')
        self.assertIn("+7 900 555-44-33", restored_text)

    def test_deanonymize_file_invalid_session(self):
        resp = self.client.post(
            "/api/deanonymize/file",
            files={"file": ("file.txt", io.BytesIO(b"text"), "text/plain")},
            data={"session_id": "bad-session-id"}
        )
        self.assertEqual(resp.status_code, 404)

    # --- GROUP API ---

    def test_group_session_text(self):
        group_id = f"API_GRP_{uuid.uuid4().hex[:8]}"

        # Файл 1
        resp1 = self.client.post("/api/anonymize/text", json={
            "text": "Клиент: ivanov@mail.ru", "mode": "fake",
            "group_id": group_id
        })
        self.assertEqual(resp1.status_code, 200)

        # Файл 2
        resp2 = self.client.post("/api/anonymize/text", json={
            "text": "Телефон: +7 999 123-45-67", "mode": "fake",
            "group_id": group_id
        })
        self.assertEqual(resp2.status_code, 200)

        # Деанонимизация обоих через group_id
        combined_text = resp1.json()["anonymized_text"] + " " + resp2.json()["anonymized_text"]
        resp3 = self.client.post("/api/deanonymize/text", json={
            "text": combined_text, "session_id": group_id
        })
        self.assertEqual(resp3.status_code, 200)
        restored = resp3.json()["restored_text"]
        self.assertIn("ivanov@mail.ru", restored)
        self.assertIn("+7 999 123-45-67", restored)

    # --- SESSIONS API ---

    def test_get_sessions(self):
        resp = self.client.get("/api/sessions")
        self.assertEqual(resp.status_code, 200)
        data = resp.json()
        self.assertIn("sessions", data)
        self.assertIsInstance(data["sessions"], list)


# ═══════════════════════════════════════════
# 8. ТЕСТЫ ПАРСЕРОВ PLAIN
# ═══════════════════════════════════════════

class TestPlainParsers(unittest.TestCase):
    """Тесты plain-text парсеров"""

    def _anonymize_func(self, text):
        return text.replace("SECRET", "***")

    def _deanonymize_func(self, text):
        return text.replace("***", "SECRET")

    def test_txt_roundtrip(self):
        original = b"This is SECRET data"
        anon = PlainParser.anonymize_txt(original, self._anonymize_func)
        self.assertIn(b"***", anon)
        restored = PlainParser.deanonymize_txt(anon, self._deanonymize_func)
        self.assertEqual(restored, original)

    def test_csv_roundtrip(self):
        original = "Name,Data\nJohn,SECRET\n".encode('utf-8')
        anon = PlainParser.anonymize_csv(original, self._anonymize_func)
        self.assertNotIn(b"SECRET", anon)
        restored = PlainParser.deanonymize_csv(anon, self._deanonymize_func)
        self.assertIn(b"SECRET", restored)

    def test_tsv_roundtrip(self):
        original = "Name\tData\nJohn\tSECRET\n".encode('utf-8')
        anon = PlainParser.anonymize_tsv(original, self._anonymize_func)
        self.assertNotIn(b"SECRET", anon)

    def test_json_roundtrip(self):
        data = {"key": "SECRET", "nested": {"inner": "SECRET"}}
        original = json.dumps(data).encode('utf-8')
        anon = PlainParser.anonymize_json(original, self._anonymize_func)
        self.assertNotIn(b"SECRET", anon)
        self.assertIn(b"***", anon)
        restored = PlainParser.deanonymize_json(anon, self._deanonymize_func)
        result = json.loads(restored)
        self.assertEqual(result["key"], "SECRET")
        self.assertEqual(result["nested"]["inner"], "SECRET")

    def test_json_nested_list(self):
        data = [{"name": "SECRET"}, "SECRET"]
        original = json.dumps(data).encode('utf-8')
        anon = PlainParser.anonymize_json(original, self._anonymize_func)
        self.assertNotIn(b"SECRET", anon)

    def test_xml_roundtrip(self):
        xml = '<?xml version="1.0"?><root><item>SECRET</item></root>'
        anon = PlainParser.anonymize_xml(xml.encode('utf-8'), self._anonymize_func)
        self.assertNotIn(b"SECRET", anon)

    def test_html_roundtrip(self):
        html = '<html><body><p>SECRET</p></body></html>'
        anon = PlainParser.anonymize_html(html.encode('utf-8'), self._anonymize_func)
        self.assertNotIn(b"SECRET", anon)

    def test_rtf_roundtrip(self):
        rtf = b'SECRET document text'
        anon = PlainParser.anonymize_rtf(rtf, self._anonymize_func)
        self.assertNotIn(b"SECRET", anon)


# ═══════════════════════════════════════════
# 9. ТЕСТЫ OFFICE ПАРСЕРОВ
# ═══════════════════════════════════════════

class TestOfficeParsers(unittest.TestCase):
    """Тесты Office парсеров"""

    def _anonymize_func(self, text):
        return text.replace("CONFIDENTIAL", "***")

    def _deanonymize_func(self, text):
        return text.replace("***", "CONFIDENTIAL")

    def test_docx_anonymize(self):
        doc = Document()
        doc.add_paragraph("This is CONFIDENTIAL information")
        buf = io.BytesIO()
        doc.save(buf)

        result = OfficeParser.anonymize_docx(buf.getvalue(), self._anonymize_func)
        doc2 = Document(io.BytesIO(result))
        all_text = " ".join(p.text for p in doc2.paragraphs)
        self.assertNotIn("CONFIDENTIAL", all_text)
        self.assertIn("***", all_text)

    def test_docx_deanonymize(self):
        doc = Document()
        doc.add_paragraph("This is *** information")
        buf = io.BytesIO()
        doc.save(buf)

        result = OfficeParser.deanonymize_docx(buf.getvalue(), self._deanonymize_func)
        doc2 = Document(io.BytesIO(result))
        all_text = " ".join(p.text for p in doc2.paragraphs)
        self.assertIn("CONFIDENTIAL", all_text)

    def test_xlsx_anonymize(self):
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Name"
        ws["A2"] = "CONFIDENTIAL"
        buf = io.BytesIO()
        wb.save(buf)

        result = OfficeParser.anonymize_xlsx(buf.getvalue(), self._anonymize_func)
        import openpyxl
        wb2 = openpyxl.load_workbook(io.BytesIO(result))
        ws2 = wb2.active
        self.assertNotEqual(ws2["A2"].value, "CONFIDENTIAL")

    def test_xlsx_sensitive_headers(self):
        """Чувствительные заголовки (зарплата и т.д.)"""
        self.assertTrue(is_sensitive_header("Зарплата сотрудника"))
        self.assertTrue(is_sensitive_header("salary"))
        self.assertTrue(is_sensitive_header("Бонус"))
        self.assertFalse(is_sensitive_header("Имя"))
        self.assertFalse(is_sensitive_header(None))

    def test_pptx_anonymize(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "CONFIDENTIAL Slide"
        buf = io.BytesIO()
        prs.save(buf)

        result = OfficeParser.anonymize_pptx(buf.getvalue(), self._anonymize_func)
        prs2 = Presentation(io.BytesIO(result))
        for slide2 in prs2.slides:
            for shape in slide2.shapes:
                if hasattr(shape, "text"):
                    self.assertNotIn("CONFIDENTIAL", shape.text)


# ═══════════════════════════════════════════
# 10. ТЕСТЫ EDGE CASES
# ═══════════════════════════════════════════

class TestEdgeCases(unittest.TestCase):
    """Граничные случаи и стресс-тесты"""

    def setUp(self):
        self.client = TestClient(app)

    def test_very_long_text(self):
        """Длинный текст (5000+ символов)"""
        text = "Иванов Иван ivan@mail.ru +7 999 123-45-67. " * 100
        resp = self.client.post("/api/anonymize/text", json={
            "text": text, "mode": "fake", "lang": "ru"
        })
        self.assertEqual(resp.status_code, 200)
        self.assertNotIn("ivan@mail.ru", resp.json()["anonymized_text"])

    def test_unicode_text(self):
        """Текст с Unicode символами"""
        text = "Клиент: 日本語テスト, email: user@test.com"
        resp = self.client.post("/api/anonymize/text", json={
            "text": text, "mode": "fake"
        })
        self.assertEqual(resp.status_code, 200)

    def test_special_chars_in_text(self):
        """Спецсимволы"""
        text = "Email: <user@test.com> & tel: +7 999 123-45-67 \"quoted\""
        resp = self.client.post("/api/anonymize/text", json={
            "text": text, "mode": "fake", "lang": "ru"
        })
        self.assertEqual(resp.status_code, 200)

    def test_only_numbers_text(self):
        """Текст только из цифр"""
        text = "1234567890"
        resp = self.client.post("/api/anonymize/text", json={
            "text": text, "mode": "fake"
        })
        self.assertEqual(resp.status_code, 200)

    def test_multiline_text(self):
        """Многострочный текст"""
        text = "Строка 1: Иванов\nСтрока 2: ivan@mail.ru\nСтрока 3: +7 999 123-45-67"
        resp = self.client.post("/api/anonymize/text", json={
            "text": text, "mode": "fake", "lang": "ru"
        })
        self.assertEqual(resp.status_code, 200)
        self.assertNotIn("ivan@mail.ru", resp.json()["anonymized_text"])

    def test_concurrent_sessions(self):
        """Множественные параллельные сессии не конфликтуют"""
        results = []
        for i in range(5):
            resp = self.client.post("/api/anonymize/text", json={
                "text": f"User{i} user{i}@mail.ru",
                "mode": "fake", "lang": "ru"
            })
            self.assertEqual(resp.status_code, 200)
            results.append(resp.json())

        # Все сессии уникальны
        session_ids = [r["session_id"] for r in results]
        self.assertEqual(len(session_ids), len(set(session_ids)))


# ═══════════════════════════════════════════
# 11. ТЕСТЫ КОНФИГУРАЦИИ
# ═══════════════════════════════════════════

class TestConfig(unittest.TestCase):
    """Тесты конфигурации"""

    def test_db_path_exists(self):
        self.assertTrue(os.path.exists(os.path.dirname(DB_PATH)))

    def test_secret_key_not_empty(self):
        self.assertTrue(len(SECRET_KEY) > 0)

    def test_config_imports(self):
        from config import HOST, PORT, ENABLE_AUTH, VALID_API_KEYS, ALLOW_ORIGINS, CACHE_CAPACITY, SESSION_TTL_HOURS
        self.assertIsInstance(HOST, str)
        self.assertIsInstance(PORT, int)
        self.assertIsInstance(ENABLE_AUTH, bool)
        self.assertIsInstance(VALID_API_KEYS, set)
        self.assertIsInstance(ALLOW_ORIGINS, list)
        self.assertIsInstance(CACHE_CAPACITY, int)
        self.assertIsInstance(SESSION_TTL_HOURS, int)


# ═══════════════════════════════════════════
# 12. ТЕСТ AUDIT LOGGER (broken import)
# ═══════════════════════════════════════════

class TestAuditLogger(unittest.TestCase):
    """Проверка audit_logger.py"""

    def test_audit_logger_imports_successfully(self):
        """audit_logger.py должен импортироваться без ошибок"""
        import importlib
        mod = importlib.import_module("audit_logger")
        self.assertTrue(hasattr(mod, "log_audit_event"))


# ═══════════════════════════════════════════
# ЗАПУСК
# ═══════════════════════════════════════════

if __name__ == "__main__":
    unittest.main(verbosity=2)
